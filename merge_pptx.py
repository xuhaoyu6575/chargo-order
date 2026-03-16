from pptx import Presentation
from pptx.util import Inches, Pt
import zipfile
import os
import re
import shutil
from pathlib import Path

def get_slide_summary(prs, idx):
    """Extract title and first 100 chars of text from a slide."""
    try:
        slide = prs.slides[idx]
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text.strip())
        text = ' '.join(parts)[:200]
        return text or f'(Slide {idx+1})'
    except:
        return f'(Slide {idx+1})'

def normalize_for_compare(text):
    """Normalize text for conflict comparison."""
    if not text:
        return ''
    t = re.sub(r'\s+', ' ', text.lower().strip())
    return t[:80]

def slides_conflict(sum1, sum2):
    """Check if two slides cover same topic (simple heuristic)."""
    n1, n2 = normalize_for_compare(sum1), normalize_for_compare(sum2)
    if not n1 or not n2:
        return False
    # Same or very similar start
    if n1 == n2:
        return True
    if n1 in n2 or n2 in n1:
        return True
    # Same first 30 chars often means same slide
    if len(n1) >= 20 and len(n2) >= 20 and n1[:30] == n2[:30]:
        return True
    return False

def main():
    base = r'f:\chargo_order'
    v1_path = os.path.join(base, 'AI研发架构图v1.pptx')
    main_path = os.path.join(base, 'AI研发架构图.pptx')
    out_path = os.path.join(base, 'AI研发架构图.pptx')
    merged_path = os.path.join(base, 'AI研发架构图_merged.pptx')

    # Load both
    prs_v1 = Presentation(v1_path)
    prs_main = Presentation(main_path)

    v1_slides = [get_slide_summary(prs_v1, i) for i in range(len(prs_v1.slides))]
    main_slides = [get_slide_summary(prs_main, i) for i in range(len(prs_main.slides))]

    # Find conflicting main slides (topic overlaps with any v1 slide)
    conflicting_idxs = []
    for mi, ms in enumerate(main_slides):
        for vs in v1_slides:
            if slides_conflict(ms, vs):
                conflicting_idxs.append(mi)
                break

    # 显式移除冲突：v1 已有封面和关键路径说明，main 中重复的需移除
    # Main slide 1 (index 0): 封面/系统研发架构图 - 与 v1 封面冲突
    # Main slide 5 (index 4): AI 赋能全链路研发流程 - 与 v1 关键路径说明冲突
    for idx in [0, 4]:
        if idx not in conflicting_idxs and idx < len(main_slides):
            conflicting_idxs.append(idx)

    non_conflicting_idxs = [i for i in range(len(main_slides)) if i not in conflicting_idxs]

    # Merge using zip manipulation (preserves formatting)
    work_dir = Path(r'f:\chargo_order\_pptx_merge')
    if work_dir.exists():
        shutil.rmtree(work_dir)
    work_dir.mkdir()

    with zipfile.ZipFile(v1_path, 'r') as z:
        z.extractall(work_dir)

    # Get slide count from v1
    ppt_dir = work_dir / 'ppt'
    slides_dir = ppt_dir / 'slides'
    v1_slide_files = sorted(slides_dir.glob('slide*.xml'), key=lambda p: int(re.search(r'slide(\d+)', p.name).group(1)))
    next_slide_num = len(v1_slide_files) + 1

    # Copy non-conflicting slides from main
    with zipfile.ZipFile(main_path, 'r') as z_main:
        for idx in non_conflicting_idxs:
            slide_name = f'slide{idx+1}.xml'
            try:
                data = z_main.read(f'ppt/slides/{slide_name}')
                # Update rId references to avoid clashes - simplified
                out_name = f'slide{next_slide_num}.xml'
                (slides_dir / out_name).write_bytes(data)
                next_slide_num += 1
            except KeyError:
                pass

    # Rebuild [Content_Types].xml, presentation.xml, _rels
    # This is complex - use python-pptx copy approach instead
    shutil.rmtree(work_dir)

    # Copy v1 to output (v1 first, unchanged), then append non-conflicting from main
    shutil.copy(v1_path, merged_path)
    prs_out = Presentation(merged_path)

    for idx in non_conflicting_idxs:
        src_slide = prs_main.slides[idx]
        layout = prs_out.slide_layouts[0]  # blank
        new_slide = prs_out.slides.add_slide(layout)
        for shape in src_slide.shapes:
            try:
                el = shape._element
                new_el = type(el)(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
            except Exception:
                pass

    prs_out.save(merged_path)
    # Overwrite main with merged result
    try:
        shutil.copy(merged_path, out_path)
        print('\nMerged to: AI研发架构图.pptx')
    except Exception:
        print('\nTarget locked. Saved as: AI研发架构图_merged.pptx')

    # Report (use ascii-safe for Windows console)
    def safe(s, n=60):
        return s.encode('ascii', 'replace').decode()[:n] if s else ''
    print('=== V1 SLIDES ===')
    for i, s in enumerate(v1_slides):
        print(f'  {i+1}. {safe(s,80)}')
    print('=== MAIN SLIDES ===')
    for i, s in enumerate(main_slides):
        print(f'  {i+1}. {safe(s,80)}')
    print('=== REMOVED (conflicts) ===')
    for i in conflicting_idxs:
        print(f'  Main #{i+1}: {safe(main_slides[i],50)}')
    print('=== DONE ===')

if __name__ == '__main__':
    main()

