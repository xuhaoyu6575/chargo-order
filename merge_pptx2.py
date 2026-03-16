from pptx import Presentation
from pptx.util import Inches, Pt
import zipfile
import os
import re
import shutil
from pathlib import Path

def get_slide_summary(prs, idx):
    """Extract title and first 100 chars of text from a slide."""
    try:
        slide = prs.slides[idx]
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text.strip())
        text = ' '.join(parts)[:200]
        return text or f'(Slide {idx+1})'
    except:
        return f'(Slide {idx+1})'

def normalize_for_compare(text):
    """Normalize text for conflict comparison."""
    if not text:
        return ''
    t = re.sub(r'\s+', ' ', text.lower().strip())
    return t[:80]

def slides_conflict(sum1, sum2):
    """Check if two slides cover same topic (simple heuristic)."""
    n1, n2 = normalize_for_compare(sum1), normalize_for_compare(sum2)
    if not n1 or not n2:
        return False
    # Same or very similar start
    if n1 == n2:
        return True
    if n1 in n2 or n2 in n1:
        return True
    # Same first 30 chars often means same slide
    if len(n1) >= 20 and len(n2) >= 20 and n1[:30] == n2[:30]:
        return True
    return False

def main():
    v1_path = r'f:\chargo_order\AI研发架构图v1.pptx'
    main_path = r'f:\chargo_order\AI研发架构图.pptx'
    out_path = r'f:\chargo_order\AI_merged.pptx'

    # Load both
    prs_v1 = Presentation(v1_path)
    prs_main = Presentation(main_path)

    v1_slides = [get_slide_summary(prs_v1, i) for i in range(len(prs_v1.slides))]
    main_slides = [get_slide_summary(prs_main, i) for i in range(len(prs_main.slides))]

    # Find conflicting main slides (topic overlaps with any v1 slide)
    conflicting_idxs = []
    for mi, ms in enumerate(main_slides):
        for vs in v1_slides:
            if slides_conflict(ms, vs):
                conflicting_idxs.append(mi)
                break

    non_conflicting_idxs = [i for i in range(len(main_slides)) if i not in conflicting_idxs]

    # Merge using zip manipulation (preserves formatting)
    work_dir = Path(r'f:\chargo_order\_pptx_merge')
    if work_dir.exists():
        shutil.rmtree(work_dir)
    work_dir.mkdir()

    with zipfile.ZipFile(v1_path, 'r') as z:
        z.extractall(work_dir)

    # Get slide count from v1
    ppt_dir = work_dir / 'ppt'
    slides_dir = ppt_dir / 'slides'
    v1_slide_files = sorted(slides_dir.glob('slide*.xml'), key=lambda p: int(re.search(r'slide(\d+)', p.name).group(1)))
    next_slide_num = len(v1_slide_files) + 1

    # Copy non-conflicting slides from main
    with zipfile.ZipFile(main_path, 'r') as z_main:
        for idx in non_conflicting_idxs:
            slide_name = f'slide{idx+1}.xml'
            try:
                data = z_main.read(f'ppt/slides/{slide_name}')
                # Update rId references to avoid clashes - simplified
                out_name = f'slide{next_slide_num}.xml'
                (slides_dir / out_name).write_bytes(data)
                next_slide_num += 1
            except KeyError:
                pass

    # Rebuild [Content_Types].xml, presentation.xml, _rels
    # This is complex - use python-pptx copy approach instead
    shutil.rmtree(work_dir)

    # Fallback: create new prs, add v1 slides by copying content, then add non-conflicting from main
    # python-pptx cannot easily copy slides. Use duplicate-and-modify approach.
    # Simpler: copy v1 to output, then append non-conflicting from main by adding new slides with content
    shutil.copy(v1_path, out_path + '.tmp')
    prs_out = Presentation(out_path + '.tmp')

    for idx in non_conflicting_idxs:
        src_slide = prs_main.slides[idx]
        layout = prs_out.slide_layouts[0]  # blank
        new_slide = prs_out.slides.add_slide(layout)
        for shape in src_slide.shapes:
            try:
                el = shape._element
                new_el = type(el)(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
            except Exception:
                pass

    prs_out.save(out_path)
    os.remove(out_path + '.tmp')

    # Report
    print('=== V1 SLIDES (AI研发架构图v1.pptx) ===')
    for i, s in enumerate(v1_slides):
        print(f'  {i+1}. {s[:80]}...' if len(s)>80 else f'  {i+1}. {s}')
    print()
    print('=== MAIN SLIDES (AI研发架构图.pptx) ===')
    for i, s in enumerate(main_slides):
        print(f'  {i+1}. {s[:80]}...' if len(s)>80 else f'  {i+1}. {s}')
    print()
    print('=== REMOVED AS CONFLICTS ===')
    for i in conflicting_idxs:
        print(f'  Main slide {i+1}: {main_slides[i][:60]}...' if len(main_slides[i])>60 else f'  Main slide {i+1}: {main_slides[i]}')
    print()
    print('=== FINAL SLIDE ORDER ===')
    for i, s in enumerate(v1_slides):
        print(f'  {i+1}. [V1] {s[:60]}...' if len(s)>60 else f'  {i+1}. [V1] {s}')
    for idx in non_conflicting_idxs:
        s = main_slides[idx]
        print(f'  {len(v1_slides)+len([x for x in non_conflicting_idxs if x<idx])+1}. [Main] {s[:60]}...' if len(s)>60 else f'  ... [Main] {s}')

if __name__ == '__main__':
    main()

