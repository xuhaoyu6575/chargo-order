import asyncio
import os
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

PPT_PATH = "f:/chargo_order/AI研发架构图_汇报.pptx"
HTML_PATH = "file:///f:/chargo_order/AI研发架构图_ppt_v1.html"
WORK_DIR = "f:/chargo_order"

async def capture_screenshots():
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": 1600, "height": 900})
        await page.goto(HTML_PATH)
        await page.wait_for_load_state("networkidle")
        
        # 0. Full slide
        await page.screenshot(path=os.path.join(WORK_DIR, "snap_full.png"), full_page=True)
        
        # 1. Title + Flow + Deliverables
        box_dlv = await page.locator('.dlv-row').bounding_box()
        clip_top = {
            "x": 0, 
            "y": 0, 
            "width": 1600, 
            "height": box_dlv['y'] + box_dlv['height'] + 20
        }
        await page.screenshot(path=os.path.join(WORK_DIR, "snap_top.png"), clip=clip_top)
        
        # 2. Frontend Column
        col1 = await page.locator('.main-grid > div:nth-child(1)').bounding_box()
        clip_col1 = {
            "x": col1['x'] - 20,
            "y": col1['y'] - 20,
            "width": col1['width'] + 40,
            "height": col1['height'] + 40
        }
        await page.screenshot(path=os.path.join(WORK_DIR, "snap_frontend.png"), clip=clip_col1)
        
        # 3. Backend Column
        col2 = await page.locator('.main-grid > div:nth-child(2)').bounding_box()
        clip_col2 = {
            "x": col2['x'] - 20,
            "y": col2['y'] - 20,
            "width": col2['width'] + 40,
            "height": col2['height'] + 40
        }
        await page.screenshot(path=os.path.join(WORK_DIR, "snap_backend.png"), clip=clip_col2)
        
        # 4. Review & Docs Column
        col3 = await page.locator('.main-grid > div:nth-child(3)').bounding_box()
        clip_col3 = {
            "x": col3['x'] - 20,
            "y": col3['y'] - 20,
            "width": col3['width'] + 40,
            "height": col3['height'] + 40
        }
        await page.screenshot(path=os.path.join(WORK_DIR, "snap_docs.png"), clip=clip_col3)
        
        # 5. Bottom Stats
        bottom = await page.locator('.bottom-bar').bounding_box()
        clip_bottom = {
            "x": 0,
            "y": bottom['y'] - 20,
            "width": 1600,
            "height": bottom['height'] + 40
        }
        await page.screenshot(path=os.path.join(WORK_DIR, "snap_bottom.png"), clip=clip_bottom)
        
        await browser.close()

def create_ppt():
    prs = Presentation()
    # Set to 16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    def add_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(5, 15, 30) # #050f1e
        
    # 1. Full slide - fill the page
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide)
    
    # Just paste image
    slide.shapes.add_picture(os.path.join(WORK_DIR, "snap_full.png"), 0, 0, width=Inches(16))
    
    # 2. Centered detail slides
    def add_centered(img_name):
        img_path = os.path.join(WORK_DIR, img_name)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        
        pic = slide.shapes.add_picture(img_path, 0, 0)
        
        # Scale to max width or height, with 5% padding
        scale_w = (prs.slide_width * 0.95) / pic.width
        scale_h = (prs.slide_height * 0.95) / pic.height
        scale = min(scale_w, scale_h)
        
        new_w = int(pic.width * scale)
        new_h = int(pic.height * scale)
        pic.width = new_w
        pic.height = new_h
        
        # Center horizontally and vertically
        pic.left = int((prs.slide_width - new_w) / 2)
        pic.top = int((prs.slide_height - new_h) / 2)

    add_centered("snap_top.png")
    add_centered("snap_frontend.png")
    add_centered("snap_backend.png")
    add_centered("snap_docs.png")
    add_centered("snap_bottom.png")
    
    prs.save(PPT_PATH)
    print(f"PPT saved to {PPT_PATH}")

if __name__ == "__main__":
    asyncio.run(capture_screenshots())
    create_ppt()
